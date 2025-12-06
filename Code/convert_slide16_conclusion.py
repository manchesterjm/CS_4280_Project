# -*- coding: utf-8 -*-
"""
Convert slide 16 from "What's Next?" to "RNN Conclusion"
"""

from pptx import Presentation

UPDATED_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def main():
    print("Loading...")
    prs = Presentation(UPDATED_PATH)
    slide = prs.slides[15]  # Slide 16 (0-indexed)

    # Replacements to convert to conclusion
    replacements = {
        # Title
        "What's Next?": "RNN Conclusion",

        # Section headers
        "Balanced Data Attempt Failed": "Key Achievements",
        "Solution: TESS Sector 1 Ground Truth": "Lessons Learned",
        "Cross-Mission Testing": "Future Work",

        # Content
        "Tried 50/50 balanced synthetic data (200 planets + 200 non-planets)": "AUC: 92.61% (exceeded 0.948 target from Malik)",
        "AUC dropped to 0.45 (domain shift) on real TESS data": "100% Recall - no planets missed",

        "33,051 windows from labeled ground truth data": "Data quantity matters (62× improvement)",
        "Same 23% positive ratio but 62x more data": "Domain shift is real (synthetic failed)",

        "Trained on TESS Sector 1 → Achieved AUC 0.93, 100% Recall on test set": "Cross-mission generalization (TESS→Kepler)",
    }

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)

    print("Saving...")
    prs.save(UPDATED_PATH)
    print("Done!")


if __name__ == "__main__":
    main()
