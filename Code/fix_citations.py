"""
Fix the citations on slides 6-9 that weren't fully replaced.
"""

from pptx import Presentation

PPTX_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def replace_full_text_in_shape(shape, old_patterns, new_text):
    """Replace entire text frame content if it matches patterns."""
    if not shape.has_text_frame:
        return False

    full_text = shape.text_frame.text
    for pattern in old_patterns:
        if pattern in full_text:
            # Clear all paragraphs and set new text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
            # Set new text in first paragraph
            if shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                if p.runs:
                    p.runs[0].text = new_text
                else:
                    p.text = new_text
            return True
    return False


def main():
    print(f"Loading: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)

    # ==========================================================================
    # SLIDE 6: Fix citation list
    # ==========================================================================
    print("Fixing Slide 6 citations...")
    slide_6 = prs.slides[5]
    for shape in slide_6.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            # Find the citation block and replace it entirely
            if "Becker" in text and "Schanche" in text and "Malik" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if shape.text_frame.paragraphs:
                    shape.text_frame.paragraphs[0].runs[0].text = """Becker, I., Pichara, K., Catelan, M., et al. (2021). Scalable end-to-end recurrent neural network for variable star classification. MNRAS, 493(2), 2981-2995.
Schanche, N., Cameron, A. C., Hébrard, G., et al. (2019). Machine-learning approaches to exoplanet transit detection. MNRAS, 483(4), 5534-5547.
Malik, A., Moster, B. P., & Obermeier, C. (2022). Exoplanet detection using machine learning. MNRAS, 513(4), 5505-5516."""

    # ==========================================================================
    # SLIDE 7: Fix Becker citation
    # ==========================================================================
    print("Fixing Slide 7 citation...")
    slide_7 = prs.slides[6]
    for shape in slide_7.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Becker" in text and ("Machine" in text or "MNRAS" in text or "learning" in text.lower()):
                if len(text) > 50:  # It's a citation, not a title
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ""
                    if shape.text_frame.paragraphs:
                        shape.text_frame.paragraphs[0].runs[0].text = "Becker, I., Pichara, K., Catelan, M., et al. (2021). Scalable end-to-end recurrent neural network for variable star classification. MNRAS, 493(2), 2981-2995."

    # ==========================================================================
    # SLIDE 8: Fix Schanche citation
    # ==========================================================================
    print("Fixing Slide 8 citation...")
    slide_8 = prs.slides[7]
    for shape in slide_8.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Schanche" in text and ("MNRAS" in text or "Harnessing" in text):
                if len(text) > 50:  # It's a citation
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ""
                    if shape.text_frame.paragraphs:
                        shape.text_frame.paragraphs[0].runs[0].text = "Schanche, N., Cameron, A. C., Hébrard, G., et al. (2019). Machine-learning approaches to exoplanet transit detection. MNRAS, 483(4), 5534-5547."

    # ==========================================================================
    # SLIDE 9: Fix Malik citation
    # ==========================================================================
    print("Fixing Slide 9 citation...")
    slide_9 = prs.slides[8]
    for shape in slide_9.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Malik" in text and ("Photometric" in text or "Monthly" in text):
                if len(text) > 50:  # It's a citation
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ""
                    if shape.text_frame.paragraphs:
                        shape.text_frame.paragraphs[0].runs[0].text = "Malik, A., Moster, B. P., & Obermeier, C. (2022). Exoplanet detection using machine learning. MNRAS, 513(4), 5505-5516."

    # Save
    print(f"\nSaving to: {PPTX_PATH}")
    prs.save(PPTX_PATH)
    print("Done!")


if __name__ == "__main__":
    main()
