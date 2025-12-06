"""
Fix all citations to APA 6th Edition format.

APA 6 format for journal articles:
Author, A. A., Author, B. B., & Author, C. C. (Year). Title of article.
    Title of Periodical, volume(issue), pages. doi:xxxxx

For 3-7 authors: List all authors in reference list
For 8+ authors: List first 6, then ... then last author
"""

from pptx import Presentation

PPTX_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"

# Correct APA 6 citations for the three final papers
BECKER_APA6 = "Becker, I., Pichara, K., Catelan, M., Protopapas, P., Aguirre, C., & Nikzat, F. (2021). Scalable end-to-end recurrent neural network for variable star classification. Monthly Notices of the Royal Astronomical Society, 493(2), 2981-2995."

SCHANCHE_APA6 = "Schanche, N., Cameron, A. C., Hébrard, G., Nielsen, L., Triaud, A. H. M. J., Almenara, J. M., ... Udry, S. (2019). Machine-learning approaches to exoplanet transit detection and candidate validation in wide-field ground-based surveys. Monthly Notices of the Royal Astronomical Society, 483(4), 5534-5547."

MALIK_APA6 = "Malik, A., Moster, B. P., & Obermeier, C. (2022). Exoplanet detection using machine learning. Monthly Notices of the Royal Astronomical Society, 513(4), 5505-5516."


def main():
    print(f"Loading: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)

    # ==========================================================================
    # SLIDE 6: Fix citation list to APA 6
    # ==========================================================================
    print("Fixing Slide 6 citations to APA 6...")
    slide_6 = prs.slides[5]
    for shape in slide_6.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Becker" in text and "Schanche" in text:
                # Clear and replace with APA 6 format
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = f"{BECKER_APA6}\n\n{SCHANCHE_APA6}\n\n{MALIK_APA6}"

    # ==========================================================================
    # SLIDE 7: Fix Becker citation to APA 6
    # ==========================================================================
    print("Fixing Slide 7 citation to APA 6...")
    slide_7 = prs.slides[6]
    for shape in slide_7.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Becker" in text and "MNRAS" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = BECKER_APA6

    # ==========================================================================
    # SLIDE 8: Fix Schanche citation to APA 6
    # ==========================================================================
    print("Fixing Slide 8 citation to APA 6...")
    slide_8 = prs.slides[7]
    for shape in slide_8.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Schanche" in text and "MNRAS" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = SCHANCHE_APA6

    # ==========================================================================
    # SLIDE 9: Fix Malik citation to APA 6
    # ==========================================================================
    print("Fixing Slide 9 citation to APA 6...")
    slide_9 = prs.slides[8]
    for shape in slide_9.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Malik" in text and "MNRAS" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = MALIK_APA6

    # Save
    print(f"\nSaving to: {PPTX_PATH}")
    prs.save(PPTX_PATH)

    print("\n" + "="*70)
    print("APA 6 CITATIONS APPLIED")
    print("="*70)
    print(f"\nSlide 6-7 (Becker):\n{BECKER_APA6}")
    print(f"\nSlide 6,8 (Schanche):\n{SCHANCHE_APA6}")
    print(f"\nSlide 6,9 (Malik):\n{MALIK_APA6}")
    print("\nDone!")


if __name__ == "__main__":
    main()
