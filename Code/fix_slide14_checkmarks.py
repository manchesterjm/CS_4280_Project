# -*- coding: utf-8 -*-
"""
Fix the checkmarks on slide 14 - Final Results should have ✓ not ❌
"""

from pptx import Presentation

UPDATED_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def main():
    print("Loading...")
    prs = Presentation(UPDATED_PATH)
    slide = prs.slides[13]

    # Replace ❌ with ✓ for the final results lines
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # Fix the final results - these should be checkmarks
                    if "AUC: 92.61%" in run.text:
                        run.text = run.text.replace("❌", "✓")
                    if "Recall: 100%" in run.text:
                        run.text = run.text.replace("❌", "✓")
                    if "F1: 0.57" in run.text:
                        run.text = run.text.replace("❌", "✓")

    print(f"Saving...")
    prs.save(UPDATED_PATH)
    print("Done!")


if __name__ == "__main__":
    main()
