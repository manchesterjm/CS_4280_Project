# -*- coding: utf-8 -*-
"""
Fix slide 14 checkmarks after merge.
"""

from pptx import Presentation

UPDATED_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def main():
    print("Loading...")
    prs = Presentation(UPDATED_PATH)
    slide = prs.slides[13]

    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # Fix the SMOTE lines - add back checkmarks
                    if run.text == "SMOTE + Optuna optimization":
                        run.text = "✓ SMOTE + Optuna optimization"
                    if run.text == "Sector 1: 33,051 windows":
                        run.text = "✓ Sector 1: 33,051 windows (62× more)"

    print("Saving...")
    prs.save(UPDATED_PATH)
    print("Done!")


if __name__ == "__main__":
    main()
