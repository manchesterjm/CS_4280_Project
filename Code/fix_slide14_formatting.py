"""
Fix Slide 14 by copying from original and doing proper text replacement
that preserves formatting.
"""

from pptx import Presentation
from copy import deepcopy

ORIGINAL_PATH = r"D:\CS_4280_Project\term_project_files\AI Midterm Presentation Final.pptx"
UPDATED_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def replace_text_preserve_format(shape, replacements):
    """Replace text while preserving formatting."""
    if not shape.has_text_frame:
        return

    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for old, new in replacements.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)


def main():
    print("Loading original presentation...")
    orig_prs = Presentation(ORIGINAL_PATH)

    print("Loading updated presentation...")
    upd_prs = Presentation(UPDATED_PATH)

    # Get original slide 14 content structure
    orig_slide_14 = orig_prs.slides[13]
    upd_slide_14 = upd_prs.slides[13]

    # First, let's see the original structure
    print("\nOriginal Slide 14 structure:")
    for i, shape in enumerate(orig_slide_14.shapes):
        if shape.has_text_frame:
            print(f"Shape {i}: {shape.text_frame.text[:100]}...")

    # Strategy: Copy text content from original, then apply replacements
    # We need to rebuild from original to get formatting back

    # The replacements we need for the final version
    replacements = {
        # Title
        "Learning from Failure": "Journey to Success",

        # Section 1: 100 Planets Only -> Initial Approach
        "100 Planets Only:": "Initial Approach (Midterm):",
        "Predicted everything as planet": "655 windows, AUC 0.76",
        "100% false positives": "Too small, overfitting risk",

        # Section 2: Solution -> Synthetic Failed
        "Solution: Add Non-Planets": "Synthetic Data Failed:",
        "100 planets + 300 non-planets": "50/50 balanced synthetic data",
        "AUC: 0.69": "AUC dropped to 0.45 (domain shift!)",
        "Real TESS: Identified TIC 307210830 (prob: 0.5959)": "Solution: TESS Sector 1 Ground Truth",

        # Section 3: But -> Solution details
        "But: Imbalanced Data Problem": "33,051 real labeled windows (62× more data)",
        "150 planets vs 505 non-planets (23% positive)": "7,686 planets + 25,365 non-planets",
        "High recall (88.67%) but low precision (38.27%)": "Final Result:",
        "Too many false positives": "AUC: 92.61% | Recall: 100%",
        "Model biased toward negatives": "F1: 0.5708 (+17% vs midterm)",
    }

    # Apply replacements to updated slide while trying to preserve structure
    for shape in upd_slide_14.shapes:
        if shape.has_text_frame:
            replace_text_preserve_format(shape, replacements)

    print("\nSaving...")
    upd_prs.save(UPDATED_PATH)
    print("Done!")


if __name__ == "__main__":
    main()
