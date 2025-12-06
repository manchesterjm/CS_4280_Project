# -*- coding: utf-8 -*-
"""
Fix Slide 14 properly by restoring from original and doing line-by-line replacements.
This preserves the original formatting.
"""

from pptx import Presentation

ORIGINAL_PATH = r"D:\CS_4280_Project\term_project_files\AI Midterm Presentation Final.pptx"
UPDATED_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def main():
    print("Loading original...")
    orig_prs = Presentation(ORIGINAL_PATH)

    print("Loading updated...")
    upd_prs = Presentation(UPDATED_PATH)

    orig_slide = orig_prs.slides[13]
    upd_slide = upd_prs.slides[13]

    # First, restore original text structure to get formatting back
    orig_shapes = [s for s in orig_slide.shapes if s.has_text_frame]
    upd_shapes = [s for s in upd_slide.shapes if s.has_text_frame]

    for orig_shape, upd_shape in zip(orig_shapes, upd_shapes):
        for orig_para, upd_para in zip(
                orig_shape.text_frame.paragraphs,
                upd_shape.text_frame.paragraphs):
            for orig_run, upd_run in zip(orig_para.runs, upd_para.runs):
                # Restore original text first
                upd_run.text = orig_run.text

    # Now do precise text replacements that fit the original structure
    # Original structure has these lines with ❌ and ✓
    replacements = {
        # Title
        "Learning from Failure": "Journey to Success",

        # Section headers (keep same structure)
        "100 Planets Only:": "Midterm (655 windows):",
        "Solution: Add Non-Planets": "SMOTE + Sector 1 Solution:",
        "But: Imbalanced Data Problem": "Final Results:",

        # Content lines - map old to new, keeping ❌ and ✓
        "Predicted everything as planet": "AUC: 0.76, small dataset",
        "100% false positives": "Synthetic failed (AUC 0.45)",

        "100 planets + 300 non-planets → AUC: 0.69": "SMOTE helped class balance",
        "Real TESS: Identified TIC 307210830 (prob: 0.5959)": "Sector 1: 33,051 windows (62× more)",

        "150 planets vs 505 non-planets (23% positive)": "AUC: 92.61% (+17% improvement)",
        "High recall (88.67%) but low precision (38.27%)": "Recall: 100% (all planets found!)",
        "Too many false positives → Model biased toward negatives": "F1: 0.57 | Precision: 39.93%",
    }

    # Apply replacements
    for shape in upd_slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)

    print(f"Saving to: {UPDATED_PATH}")
    upd_prs.save(UPDATED_PATH)
    print("Done!")


if __name__ == "__main__":
    main()
