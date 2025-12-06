# -*- coding: utf-8 -*-
"""
Merge slide 15 content into slide 14 and delete slide 15.
"""

from pptx import Presentation

UPDATED_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def main():
    print("Loading...")
    prs = Presentation(UPDATED_PATH)

    # Update slide 14 to include Optuna info
    slide_14 = prs.slides[13]

    # Find the content shape and update it
    for shape in slide_14.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # Add Optuna to the SMOTE section
                    if "SMOTE helped class balance" in run.text:
                        run.text = "SMOTE + Optuna optimization"
                    if "Sector 1: 33,051 windows" in run.text:
                        run.text = "Sector 1: 33,051 windows"

    # Delete slide 15 (index 14)
    print("Removing slide 15...")
    slide_id = prs.slides._sldIdLst[14].rId
    prs.part.drop_rel(slide_id)
    del prs.slides._sldIdLst[14]

    print(f"Saving...")
    prs.save(UPDATED_PATH)
    print("Done! Slide 15 removed, content merged into slide 14.")


if __name__ == "__main__":
    main()
