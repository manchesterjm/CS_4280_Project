"""
Restore Slide 14 from original and apply updates preserving formatting.
"""

from pptx import Presentation

ORIGINAL_PATH = r"D:\CS_4280_Project\term_project_files\AI Midterm Presentation Final.pptx"
UPDATED_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def main():
    print("Loading original...")
    orig_prs = Presentation(ORIGINAL_PATH)

    print("Loading updated...")
    upd_prs = Presentation(UPDATED_PATH)

    orig_slide = orig_prs.slides[13]
    upd_slide = upd_prs.slides[13]

    # Copy text content from original shapes to updated shapes
    # Match by shape index
    orig_shapes = [s for s in orig_slide.shapes if s.has_text_frame]
    upd_shapes = [s for s in upd_slide.shapes if s.has_text_frame]

    print(f"\nOriginal has {len(orig_shapes)} text shapes")
    print(f"Updated has {len(upd_shapes)} text shapes")

    # For each shape in updated, copy the paragraph/run structure from original
    # then apply our text changes
    for i, (orig_shape, upd_shape) in enumerate(zip(orig_shapes, upd_shapes)):
        print(f"\nShape {i}:")
        orig_text = orig_shape.text_frame.text[:80].encode('ascii', 'replace').decode('ascii')
        print(f"  Original: {orig_text}...")

        # Copy each paragraph's runs from original
        for p_idx, (orig_para, upd_para) in enumerate(zip(
                orig_shape.text_frame.paragraphs,
                upd_shape.text_frame.paragraphs)):

            for r_idx, (orig_run, upd_run) in enumerate(zip(
                    orig_para.runs,
                    upd_para.runs)):
                # Copy the original text first (restores structure)
                upd_run.text = orig_run.text

    # Now apply the text replacements
    replacements = {
        # Title
        "Learning from Failure": "Journey to Success",

        # Content updates
        "100 Planets Only:": "Initial Approach (Midterm):",
        "Predicted everything as planet": "655 windows, AUC 0.76",
        "100% false positives": "Too small, overfitting risk",

        "Solution: Add Non-Planets": "Synthetic Data Failed:",
        "100 planets + 300 non-planets": "50/50 balanced synthetic",
        "AUC: 0.69": "AUC: 0.45 (domain shift!)",
        "Real TESS: Identified TIC 307210830 (prob: 0.5959)": "",

        "But: Imbalanced Data Problem": "Solution: Sector 1 Ground Truth",
        "150 planets vs 505 non-planets (23% positive)": "33,051 windows (62× more data)",
        "High recall (88.67%) but low precision (38.27%)": "7,686 planets + 25,365 non-planets",
        "Too many false positives": "Final: AUC 92.61%, Recall 100%",
        "Model biased toward negatives": "F1: 0.5708 (+17% improvement)",
    }

    for shape in upd_slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)

    print("\nSaving...")
    upd_prs.save(UPDATED_PATH)
    print("Done!")

    # Verify
    print("\n" + "="*70)
    print("SLIDE 14 AFTER UPDATE:")
    print("="*70)
    verify_prs = Presentation(UPDATED_PATH)
    for shape in verify_prs.slides[13].shapes:
        if shape.has_text_frame:
            print(shape.text_frame.text)
            print("---")


if __name__ == "__main__":
    main()
