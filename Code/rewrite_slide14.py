# -*- coding: utf-8 -*-
"""
Rewrite Slide 14 with proper journey including SMOTE.
"""

from pptx import Presentation

UPDATED_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def main():
    print("Loading presentation...")
    prs = Presentation(UPDATED_PATH)

    slide = prs.slides[13]

    # New content for slide 14 - Journey to Success
    new_title = "Journey to Success"

    new_content = """Midterm (655 windows):
❌ Small dataset → overfitting
❌ AUC: 0.76, limited generalization

Synthetic Data Attempt:
❌ 50/50 balanced synthetic data
❌ AUC dropped to 0.45 (domain shift!)

SMOTE Augmentation:
✓ Oversampled minority class
✓ Helped with class balance

Final Solution: Sector 1 Ground Truth
✓ 33,051 real labeled windows (62× more)
✓ 7,686 planets + 25,365 non-planets
✓ AUC: 92.61% | Recall: 100% | F1: 0.57"""

    # Update the shapes
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text

            # Update title
            if "Journey" in text or "Learning" in text or "Success" in text:
                if len(text) < 50:  # It's the title
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = new_title

            # Update main content
            elif "Midterm" in text or "Initial" in text or "100 Planets" in text or "Approach" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                # Set new content in first run
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = new_content

    print(f"Saving to: {UPDATED_PATH}")
    prs.save(UPDATED_PATH)
    print("Done!")


if __name__ == "__main__":
    main()
