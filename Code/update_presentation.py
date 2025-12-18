"""
Update Final Presentation slides with latest RNN results.

Updates slides 10-17 (RNN section by Josh Manchester) with:
- Final model results: AUC 0.9261, 100% Recall
- Dataset: 33,051 windows (26,472 train / 6,579 test)
- Architecture: 4-layer BiLSTM, 192 hidden, 7 clusters, 3.07M params
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os

# Paths
PPTX_PATH = r"D:\CS_4280_Project\term_project_files\AI Midterm Presentation Final.pptx"
OUTPUT_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def analyze_presentation(prs):
    """Print structure of all slides to understand what we're working with."""
    print(f"\n{'='*70}")
    print(f"Presentation has {len(prs.slides)} slides")
    print(f"{'='*70}\n")

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n--- Slide {slide_idx + 1} ---")
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text_preview = shape.text_frame.text[:100].replace('\n', ' ')
                # Handle Unicode characters that can't be printed
                text_preview = text_preview.encode('ascii', 'replace').decode('ascii')
                print(f"  Shape {shape_idx}: {shape.shape_type} - '{text_preview}...'")


def find_and_replace_text(shape, old_text, new_text):
    """Find and replace text in a shape, preserving formatting."""
    if not shape.has_text_frame:
        return False

    replaced = False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                replaced = True
    return replaced


def update_slide_text(slide, replacements):
    """Apply multiple text replacements to a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for old_text, new_text in replacements.items():
                find_and_replace_text(shape, old_text, new_text)


def update_rnn_slides(prs):
    """Update RNN section slides (10-17) with final results."""

    # Define replacements for each slide
    # Note: Slide indices are 0-based, so slide 10 is index 9

    # Slide 11 (index 10) - Methodology Overview
    slide_11_replacements = {
        "655 windows": "33,051 windows",
        "150 planets, 505 non-planets": "7,686 planets, 25,365 non-planets",
        "(150 planets": "(7,686 planets",
        "505 non-planets)": "25,365 non-planets)",
        "5 clusters": "7 clusters",
        "~2.1M parameters": "3.07M parameters",
        "256 hidden": "192 hidden",
        "3 layers": "4 layers",
        "32-dimensional": "64-dimensional",
        "32-dim": "64-dim",
        "3-layer BiLSTM": "4-layer BiLSTM",
    }

    # Slide 12 (index 11) - Architecture
    slide_12_replacements = {
        "~2.1M": "3.07M",
        "2.1M": "3.07M",
        "256 hidden units": "192 hidden units",
        "3 layers": "4 layers",
        "5 clusters": "7 clusters",
        "32-dim": "64-dim",
        "Dropout 0.4": "Dropout 0.334",
        "0.4": "0.334",
    }

    # Slide 13 (index 12) - Results
    slide_13_replacements = {
        # AUC updates
        "75.72": "92.61",
        "0.7572": "0.9261",
        "AUC: 0.76": "AUC: 0.93",
        "AUC: 76": "AUC: 93",
        "AUC 0.76": "AUC 0.93",
        "AUC 76": "AUC 93",
        # Recall updates
        "88.67": "100",
        "0.8867": "1.0",
        "Recall: 89": "Recall: 100",
        "Recall: 88": "Recall: 100",
        # F1 updates
        "0.4550": "0.5708",
        "45.50": "57.08",
        "F1: 0.46": "F1: 0.57",
        "F1: 46": "F1: 57",
        # Precision updates
        "30.94": "39.93",
        "0.3094": "0.3993",
        "38.27": "39.93",
        "Precision: 38.27%": "Precision: 39.93%",
        "100 confirmed exoplanet systems": "6,579 test windows (TESS Sector 1)",
        # Accuracy updates
        "52%": "83.26%",
        "52.0": "83.26",
    }

    # Slide 14 (index 13) - Confusion Matrix
    slide_14_replacements = {
        # Old values to new values
        "TP=5": "TP=2453",
        "FP=8": "FP=3687",
        "TN=43": "TN=1858",
        "FN=45": "FN=0",
    }

    # Slide 15 (index 14) - Real TESS Testing
    slide_15_replacements = {
        "16/300": "2453/2453",
        "5.3%": "100%",
        "100 confirmed": "6,579 test windows",
    }

    # Slide 16 (index 15) - Improvement
    slide_16_replacements = {
        "+9.0%": "+18.9%",
        "0.6947": "0.7572",  # baseline
        "baseline": "midterm (0.76)",
        "AUC 0.69": "AUC 0.93",
    }

    # Apply replacements to each slide
    print("\nUpdating slides...")

    # Process all slides and apply relevant replacements
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        # Get all text from slide for analysis
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text += shape.text_frame.text + " "

        # Apply replacements based on slide content patterns
        replacements = {}

        # Common RNN result replacements (apply to all RNN slides)
        if any(keyword in slide_text.lower() for keyword in ['auc', 'recall', 'precision', 'f1', 'bilstm', 'lstm', 'cluster']):
            replacements.update(slide_13_replacements)
            replacements.update(slide_11_replacements)
            replacements.update(slide_12_replacements)

        if replacements:
            update_slide_text(slide, replacements)
            print(f"  Updated slide {slide_num}")


def main():
    """Main function to update presentation."""
    print(f"Loading presentation from: {PPTX_PATH}")

    if not os.path.exists(PPTX_PATH):
        print(f"ERROR: File not found: {PPTX_PATH}")
        return

    prs = Presentation(PPTX_PATH)

    # Analyze first to see structure
    analyze_presentation(prs)

    # Update RNN slides
    update_rnn_slides(prs)

    # Save updated presentation
    print(f"\nSaving updated presentation to: {OUTPUT_PATH}")
    prs.save(OUTPUT_PATH)
    print("Done!")


if __name__ == "__main__":
    main()
