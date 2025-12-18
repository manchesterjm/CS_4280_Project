"""
Update Final Presentation - ALL RNN slides (6-17, skip 16 demo).

Final model results (December 6, 2025):
- Dataset: 33,051 windows (26,472 train / 6,579 test)
- AUC: 0.9261, Recall: 100%, Precision: 39.93%, F1: 0.5708
- Architecture: 4-layer BiLSTM, 192 hidden, 7 clusters, 64-dim embed
- Parameters: 3,068,801 (3.07M)
"""

from pptx import Presentation
import os

PPTX_PATH = r"D:\CS_4280_Project\term_project_files\AI Midterm Presentation Final.pptx"
OUTPUT_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def find_and_replace_in_shape(shape, replacements):
    """Replace text in shape, preserving formatting."""
    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for old_text, new_text in replacements.items():
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)


def update_slide(slide, replacements):
    """Apply all replacements to a slide."""
    for shape in slide.shapes:
        find_and_replace_in_shape(shape, replacements)


def main():
    """Update all RNN slides."""
    print(f"Loading: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)

    # ==========================================================================
    # SLIDE 6: New RNN Related Works (title slide - no changes needed)
    # ==========================================================================
    # Just citations, no results to update

    # ==========================================================================
    # SLIDE 7: Machine Learning for Cluster Analysis (Speiser 2020)
    # No changes needed - this is about the paper, not our results
    # ==========================================================================

    # ==========================================================================
    # SLIDE 8: LSTM for Time Series Patterns (Vu 2024)
    # No changes needed - this is about the paper, not our results
    # ==========================================================================

    # ==========================================================================
    # SLIDE 9: LSTM for Astronomical Photometry (Ding 2024)
    # No changes needed - this is about the paper, not our results
    # ==========================================================================

    # ==========================================================================
    # SLIDE 10: Why BiLSTM + Clustering?
    # ==========================================================================
    # No changes needed - this explains the motivation

    # ==========================================================================
    # SLIDE 11: Methodology
    # ==========================================================================
    print("Updating Slide 11 (Methodology)...")
    slide_11 = prs.slides[10]
    replacements_11 = {
        # Dataset
        "655 windows": "33,051 windows",
        "150 planets, 505 non-planets": "7,686 planets, 25,365 non-planets",
        # Architecture
        "5 clusters": "7 clusters",
        "4-layer BiLSTM (256 hidden": "4-layer BiLSTM (192 hidden",
        "3-layer BiLSTM": "4-layer BiLSTM",
        "(256 hidden": "(192 hidden",
        "256 hidden": "192 hidden",
        "32-dim": "64-dim",
        "(32-dim)": "(64-dim)",
    }
    update_slide(slide_11, replacements_11)

    # ==========================================================================
    # SLIDE 12: BiLSTM Architecture
    # ==========================================================================
    print("Updating Slide 12 (Architecture)...")
    slide_12 = prs.slides[11]
    replacements_12 = {
        "~2.1M parameters": "3.07M parameters",
        "2.1M parameters": "3.07M parameters",
        "~2.1M": "3.07M",
    }
    update_slide(slide_12, replacements_12)

    # ==========================================================================
    # SLIDE 13: Results
    # ==========================================================================
    print("Updating Slide 13 (Results)...")
    slide_13 = prs.slides[12]
    replacements_13 = {
        # Main metrics
        "AUC: 75.72%": "AUC: 92.61%",
        "75.72%": "92.61%",
        "Recall: 88.67%": "Recall: 100%",
        "88.67%": "100%",
        "Precision: 38.27%": "Precision: 39.93%",
        "38.27%": "39.93%",
        # Test set description
        "100 confirmed exoplanet systems": "6,579 test windows (TESS Sector 1)",
        "Tested on 100 confirmed": "Tested on 6,579 test windows",
    }
    update_slide(slide_13, replacements_13)

    # ==========================================================================
    # SLIDE 14: Learning from Failure
    # ==========================================================================
    print("Updating Slide 14 (Learning from Failure)...")
    slide_14 = prs.slides[13]
    replacements_14 = {
        # Update the "But" section with final results
        "150 planets vs 505 non-planets (23% positive)": "7,686 planets vs 25,365 non-planets (23% positive)",
        "High recall (88.67%) but low precision (38.27%)": "High recall (100%) but precision (39.93%)",
        "(88.67%)": "(100%)",
        "(38.27%)": "(39.93%)",
        "88.67%": "100%",
        "38.27%": "39.93%",
        # Update solution section
        "100 planets + 300 non-planets": "7,686 planets + 25,365 non-planets",
        "AUC: 0.69": "AUC: 0.93",
    }
    update_slide(slide_14, replacements_14)

    # ==========================================================================
    # SLIDE 15: Optuna Optimization
    # ==========================================================================
    print("Updating Slide 15 (Optuna Optimization)...")
    slide_15 = prs.slides[14]
    replacements_15 = {
        # Final result - update to show full progression
        "Result: AUC 0.69": "Result: Midterm 0.76",
        "0.76 (+9%)": "0.93 (final) (+22% total)",
    }
    update_slide(slide_15, replacements_15)

    # ==========================================================================
    # SLIDE 16: Demo - SKIP (user requested)
    # ==========================================================================

    # ==========================================================================
    # SLIDE 17: What's Next? -> Now becomes "Final Results & Lessons Learned"
    # ==========================================================================
    print("Updating Slide 17 (What's Next / Final Results)...")
    slide_17 = prs.slides[16]
    replacements_17 = {
        # Update the failed experiment section
        "AUC dropped to 0.45": "AUC dropped to 0.45 (domain shift)",
        # Update the solution section
        "Solution: Hybrid Training": "Solution: TESS Sector 1 Ground Truth",
        "Mix real TESS + synthetic (90/10 ratio)": "33,051 windows from labeled ground truth data",
        "Better balance than 23% but maintains domain fidelity": "Same 23% positive ratio but 62x more data",
        # Update cross-mission section
        "Train on TESS": "Trained on TESS Sector 1",
        "Test on Kepler to verify generalization": "Achieved AUC 0.93, 100% Recall on test set",
    }
    update_slide(slide_17, replacements_17)

    # Save
    print(f"\nSaving to: {OUTPUT_PATH}")
    prs.save(OUTPUT_PATH)
    print("Done!")

    # Print summary of changes
    print("\n" + "="*70)
    print("SUMMARY OF UPDATES")
    print("="*70)
    print("""
Slide 6-10: Related Works & Motivation - No changes (paper content)

Slide 11 (Methodology):
  - 655 windows -> 33,051 windows
  - 150/505 -> 7,686/25,365
  - 5 clusters -> 7 clusters
  - 256 hidden -> 192 hidden
  - 32-dim -> 64-dim embed

Slide 12 (Architecture):
  - ~2.1M -> 3.07M parameters

Slide 13 (Results):
  - AUC: 75.72% -> 92.61%
  - Recall: 88.67% -> 100%
  - Precision: 38.27% -> 39.93%
  - "100 confirmed exoplanet systems" -> "6,579 test windows (TESS Sector 1)"

Slide 14 (Learning from Failure):
  - Updated class counts
  - Recall/Precision updated to final values
  - AUC: 0.69 -> 0.93

Slide 15 (Optuna Optimization):
  - Updated to show progression: 0.69 -> 0.76 -> 0.93

Slide 16: SKIPPED (Demo)

Slide 17 (What's Next):
  - Updated to reflect final solution (Sector 1 ground truth)
  - Final results: AUC 0.93, 100% Recall
""")


if __name__ == "__main__":
    main()
