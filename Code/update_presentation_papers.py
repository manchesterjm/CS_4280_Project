"""
Update Presentation Slides 6-9 with the correct FINAL three papers.

Midterm papers (OLD - to be replaced):
- Speiser 2020 - Machine Learning for Cluster Analysis
- Vu 2024 - LSTM for Time Series Patterns
- Ding 2024 - LSTM for Astronomical Photometry

Final papers (NEW - correct):
- Becker et al. (2021) - CNN-LSTM for Variable Star Classification
- Schanche et al. (2019) - ML for Ground-Based Transit Surveys
- Malik et al. (2022) - Feature-Based Detection with Gradient Boosting
"""

from pptx import Presentation
from pptx.util import Pt
import os

PPTX_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"
OUTPUT_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def find_and_replace_in_shape(shape, replacements):
    """Replace text in shape, preserving formatting."""
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for old_text, new_text in replacements.items():
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)


def update_slide(slide, replacements):
    """Apply all replacements to a slide."""
    for shape in slide.shapes:
        find_and_replace_in_shape(shape, replacements)


def main():
    """Update slides 6-9 with correct final papers."""
    print(f"Loading: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)

    # ==========================================================================
    # SLIDE 6: Title slide for new papers section
    # ==========================================================================
    print("Updating Slide 6 (Paper Citations)...")
    slide_6 = prs.slides[5]
    replacements_6 = {
        # Replace old paper citations with new ones
        "Speiser, A., Müller, L. R., Matti, U., Obholzer, N. D., Legant, W. R., Kreshuk, A., ... & Hufnagel, L. (2020). Machine learning for cluster analysis of localization microscopy data. Nature Communications, 11(1), 1493. (H5 index: 399)":
            "Becker, I., Pichara, K., Catelan, M., Protopapas, P., Aguirre, C., & Nikzat, F. (2021). Scalable end-to-end recurrent neural network for variable star classification. Monthly Notices of the Royal Astronomical Society, 493(2), 2981-2995.",

        "Vu, M. T., Vo, N. D., Ngo, T. D., Pham, T. D., Huynh, T. T. M., Nguyen, N. T., ... & Ly, H. B. (2024). Harnessing LSTM and XGBoost algorithms for storm prediction. Scientific Reports, 14(1), 11516. (H5 index: 234)":
            "Schanche, N., Cameron, A. C., Hébrard, G., Nielsen, L., Triaud, A. H., Almenara, J. M., ... & Udry, S. (2019). Machine-learning approaches to exoplanet transit detection and candidate validation in wide-field ground-based surveys. Monthly Notices of the Royal Astronomical Society, 483(4), 5534-5547.",

        "Ding, Y., Ji, K., Xiao, M., Zheng, X., Chen, Y., Liang, J., ... & Qi, Y. (2024). Photometric redshift estimation for CSST survey with LSTM neural networks. Monthly Notices of the Royal Astronomical Society, 535(2), 1844-1858. (H5 index: 151)":
            "Malik, A., Moster, B. P., & Obermeier, C. (2022). Exoplanet detection using machine learning. Monthly Notices of the Royal Astronomical Society, 513(4), 5505-5516.",

        # Shorter versions in case formatting differs
        "Speiser, A.,": "Becker, I.,",
        "Vu, M. T.,": "Schanche, N.,",
        "Ding, Y.,": "Malik, A.,",
    }
    update_slide(slide_6, replacements_6)

    # ==========================================================================
    # SLIDE 7: First paper - was Speiser (Clustering), now Becker (CNN-LSTM)
    # ==========================================================================
    print("Updating Slide 7 (Paper 1: Becker - CNN-LSTM)...")
    slide_7 = prs.slides[6]
    replacements_7 = {
        # Title
        "Machine Learning for Cluster Analysis": "CNN-LSTM for Variable Star Classification",

        # Key Innovation
        "Supervised clustering combined with machine learning for analyzing millions of data points":
            "Hybrid CNN-LSTM architecture achieves 92-95% accuracy classifying variable star light curves",

        # Key Takeaway
        "K-means clustering before ML training improves both accuracy and computational speed on large-scale datasets":
            "Raw time-series input works well; weighted loss handles severe class imbalance (10-20x) effectively",

        # Architecture
        "K-means clustering on extracted features":
            "1D CNN layers (kernel 3-7) for local feature extraction",
        "Supervised learning on cluster assignments":
            "LSTM layers capture long-range temporal dependencies",
        "Scalable to millions of data points":
            "Hybrid outperforms single-architecture by 5-8%",

        # Results
        "Faster training convergence":
            "92-95% accuracy on OGLE and CRTS surveys",
        "Improved classification accuracy":
            "Cross-survey generalization validated",
        "Handles large-scale datasets efficiently":
            "Class imbalance handled via stratified sampling + weighted loss",

        # Citation
        "Speiser, A., Müller, L. R., Matti, U., Obholzer, N. D., Legant, W. R., Kreshuk, A., ... & Hufnagel, L. (2020). Machine\nlearning for cluster analysis of localization microscopy data. Nature Communications, 11(1), 1493.":
            "Becker, I., Pichara, K., Catelan, M., et al. (2021). Scalable end-to-end recurrent neural network for variable star classification. MNRAS, 493(2), 2981-2995.",

        # Shorter matches
        "Speiser, A.,": "Becker, I.,",
        "Nature Communications": "MNRAS",
    }
    update_slide(slide_7, replacements_7)

    # ==========================================================================
    # SLIDE 8: Second paper - was Vu (LSTM storms), now Schanche (Ground-Based)
    # ==========================================================================
    print("Updating Slide 8 (Paper 2: Schanche - Ground-Based Surveys)...")
    slide_8 = prs.slides[7]
    replacements_8 = {
        # Title
        "LSTM for Time Series Patterns": "ML for Ground-Based Transit Surveys",

        # Key Innovation
        "LSTM networks capture long-term dependencies in noisy environmental time series data":
            "Ensemble classifiers achieve ~90% planet identification on noisy ground-based WASP data",

        # Key Takeaway
        "LSTM memory cells handle temporal dependencies that traditional methods miss in sequential data with high noise":
            "If ML works on noisy ground data, BiLSTM should work even better on cleaner space-based TESS photometry",

        # Architecture
        "Stacked LSTM layers for temporal pattern recognition":
            "Random Forest + CNN ensemble classifiers",
        "Handles irregular sampling and data gaps":
            "Trained on planets, EBs, variable stars, non-periodic",
        "Dropout regularization for noisy data":
            "Ensemble voting reduces individual classifier errors",

        # Results
        "Superior performance on noisy time series":
            "~90% correct planet identification",
        "Captures long-term dependencies effectively":
            "Separates planets from eclipsing binary false positives",
        "Generalizes to unseen weather patterns":
            "Transit depth & duration are key discriminating features",

        # Citation
        "Vu, M. T., Vo, N. D., Ngo, T. D., Pham, T. D., Huynh, T. T. M., Nguyen, N. T., ... & Ly, H. B. (2024). Harnessing\nLSTM and XGBoost algorithms for storm prediction. Scientific Reports, 14(1), 11516.":
            "Schanche, N., Cameron, A. C., Hébrard, G., et al. (2019). Machine-learning approaches to exoplanet transit detection. MNRAS, 483(4), 5534-5547.",

        # Shorter matches
        "Vu, M. T.,": "Schanche, N.,",
        "Scientific Reports": "MNRAS",
    }
    update_slide(slide_8, replacements_8)

    # ==========================================================================
    # SLIDE 9: Third paper - was Ding (LSTM photometry), now Malik (Gradient Boosting)
    # ==========================================================================
    print("Updating Slide 9 (Paper 3: Malik - Gradient Boosting)...")
    slide_9 = prs.slides[8]
    replacements_9 = {
        # Title
        "LSTM for Astronomical Photometry": "Feature-Based Detection with Gradient Boosting",

        # Key Innovation
        "First application of LSTM to large-scale astronomical photometric flux measurements":
            "789 TSFRESH statistical features + LightGBM achieves AUC 0.948, Recall 0.96 on Kepler",

        # Key Takeaway
        "LSTM reduced outliers by 33% compared to traditional neural networks on real astronomical survey data":
            "Feature-based methods set a high bar (AUC 0.948) that neural networks must exceed to justify complexity",

        # Architecture
        "LSTM for photometric time series":
            "TSFRESH extracts 789 time series features automatically",
        "Handles large survey datasets":
            "LightGBM gradient boosting for classification",
        "Optimized for astronomical flux measurements":
            "No neural network required - pure feature engineering",

        # Results
        "33% fewer outliers than MLPs":
            "AUC: 0.948 (ranks planets above non-planets 95% of time)",
        "Better handling of temporal photometric data":
            "Recall: 0.96 (detects 96% of real planets)",
        "Validated on real survey observations":
            "Establishes performance target for neural approaches",

        # Citation
        "Ding, Y., Ji, K., Xiao, M., Zheng, X., Chen, Y., Liang, J., ... & Qi, Y. (2024). Photometric redshift estimation for CSST\nsurvey with LSTM neural networks. Monthly Notices of the Royal Astronomical Society, 535(2), 1844-1858.":
            "Malik, A., Moster, B. P., & Obermeier, C. (2022). Exoplanet detection using machine learning. MNRAS, 513(4), 5505-5516.",

        # Shorter matches
        "Ding, Y.,": "Malik, A.,",
    }
    update_slide(slide_9, replacements_9)

    # ==========================================================================
    # SLIDE 10: Why BiLSTM + Clustering? - Update paper references
    # ==========================================================================
    print("Updating Slide 10 (Motivation - paper references)...")
    slide_10 = prs.slides[9]
    replacements_10 = {
        # Update paper references in motivation
        "Clustering improves ML on large datasets (Speiser 2020)":
            "CNN-LSTM hybrid outperforms single architectures (Becker 2021)",
        "LSTM captures long-term dependencies (Vu 2024)":
            "Ensemble/clustering reduces classification errors (Schanche 2019)",
        "LSTM reduces outliers in astronomy (Ding 2024)":
            "Feature-based methods set AUC 0.948 target (Malik 2022)",

        # Also handle variants without year
        "(Speiser 2020)": "(Becker 2021)",
        "(Vu 2024)": "(Schanche 2019)",
        "(Ding 2024)": "(Malik 2022)",
    }
    update_slide(slide_10, replacements_10)

    # Save
    print(f"\nSaving to: {OUTPUT_PATH}")
    prs.save(OUTPUT_PATH)
    print("Done!")

    # Print summary
    print("\n" + "="*70)
    print("PAPER UPDATES SUMMARY")
    print("="*70)
    print("""
Slide 6 (Citations):
  OLD: Speiser (2020), Vu (2024), Ding (2024)
  NEW: Becker (2021), Schanche (2019), Malik (2022)

Slide 7:
  OLD: "Machine Learning for Cluster Analysis" (Speiser)
  NEW: "CNN-LSTM for Variable Star Classification" (Becker)
       - 92-95% accuracy on OGLE/CRTS surveys
       - Hybrid architecture outperforms single by 5-8%
       - Weighted loss handles class imbalance

Slide 8:
  OLD: "LSTM for Time Series Patterns" (Vu)
  NEW: "ML for Ground-Based Transit Surveys" (Schanche)
       - ~90% planet identification on WASP data
       - Ensemble classifiers reduce errors
       - Depth & duration are key features

Slide 9:
  OLD: "LSTM for Astronomical Photometry" (Ding)
  NEW: "Feature-Based Detection with Gradient Boosting" (Malik)
       - AUC 0.948, Recall 0.96 on Kepler
       - 789 TSFRESH features + LightGBM
       - Sets performance target for neural approaches

Slide 10 (Why BiLSTM + Clustering):
  - Updated all paper references to match new papers
""")


if __name__ == "__main__":
    main()
