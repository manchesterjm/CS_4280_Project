"""
Update Slide 14 with final project journey - from failure to success.

The slide should tell the story:
1. Initial failure (655 windows, AUC 0.69)
2. Synthetic data attempt failed (domain shift)
3. Solution: TESS Sector 1 Ground Truth (33,051 windows)
4. Final success: AUC 0.9261, 100% Recall
"""

from pptx import Presentation
from pptx.util import Pt

PPTX_PATH = r"D:\CS_4280_Project\term_project_files\AI Final Presentation - Updated.pptx"


def main():
    print(f"Loading: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)

    slide_14 = prs.slides[13]

    # Find and update the content shapes
    for shape in slide_14.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text

            # Update the title if it's "Learning from Failure"
            if "Learning from Failure" in text:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Learning from Failure" in run.text:
                            run.text = "Journey to Success"

            # Update the main content block
            if "100 Planets Only" in text or "Imbalanced Data" in text:
                # Clear all content
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""

                # Set new content
                new_content = """Initial Approach (Midterm):
• 655 windows, AUC 0.76
• Too small, overfitting risk

Synthetic Data Failed:
• 50/50 balanced synthetic data
• AUC dropped to 0.45 (domain shift!)

Solution: TESS Sector 1 Ground Truth
• 33,051 real labeled windows (62× more data)
• 7,686 planets + 25,365 non-planets

Final Result:
• AUC: 92.61% (+17% improvement)
• Recall: 100% (all planets detected!)
• F1: 0.5708"""

                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    shape.text_frame.paragraphs[0].runs[0].text = new_content

    print(f"Saving to: {PPTX_PATH}")
    prs.save(PPTX_PATH)
    print("Done! Slide 14 updated.")


if __name__ == "__main__":
    main()
