"""
Convert MIDTERM_SLIDES.md to PowerPoint matching UCCS template

Usage:
    python convert_slides_to_pptx_v2.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re

def create_blank_presentation():
    """Create presentation with UCCS-style template"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs

def add_footer(slide):
    """Add black footer bar (UCCS template style)"""
    # Add black rectangle at bottom
    left = Inches(0)
    top = Inches(7.0)  # Bottom 0.5 inches
    width = Inches(10)
    height = Inches(0.5)

    shape = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.line.color.rgb = RGBColor(0, 0, 0)

    # Add UCCS text (left side)
    left_text = slide.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(3), Inches(0.3))
    left_text.text = "UCCS - University of Colorado Colorado Springs"
    left_text.text_frame.paragraphs[0].font.size = Pt(9)
    left_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Add CU text (right side)
    right_text = slide.shapes.add_textbox(Inches(6.5), Inches(7.1), Inches(3.2), Inches(0.3))
    right_text.text = "University of Colorado"
    right_text.text_frame.paragraphs[0].font.size = Pt(9)
    right_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    right_text.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

def add_title_slide(prs, title, authors):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Add authors
    author_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(1.5))
    author_frame = author_box.text_frame
    author_frame.word_wrap = True

    for author in authors:
        p = author_frame.add_paragraph()
        p.text = author
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.alignment = PP_ALIGN.CENTER

    add_footer(slide)
    return slide

def add_content_slide(prs, title, bullets, max_bullets=10):
    """Add content slide with title and bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    # Add content area
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.2))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.clear()

    bullet_count = 0
    for bullet in bullets[:max_bullets]:
        if bullet_count >= max_bullets:
            break

        # Parse bullet level
        level = 0
        text = bullet.strip()

        if text.startswith('- '):
            text = text[2:].strip()
            level = 0
        elif text.startswith('  - '):
            text = text[4:].strip()
            level = 1
        elif text.startswith('o '):
            text = text[2:].strip()
            level = 1

        # Skip empty or header lines
        if not text or text.startswith('#') or text == '---':
            continue

        # Clean text
        text = clean_markdown(text)

        # Detect APA citations (lines with author names, years, DOI/URL)
        is_citation = ('https://doi.org' in text or
                      'https://arxiv.org' in text or
                      re.search(r'\(\d{4}\)', text) is not None)

        # Add paragraph
        if bullet_count == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()

        p.text = text
        p.level = level if not is_citation else 0

        # Smaller font for citations to fit more text
        if is_citation:
            p.font.size = Pt(14)
            p.space_after = Pt(12)  # More space between citations
        else:
            p.font.size = Pt(18) if level == 0 else Pt(16)
            p.space_after = Pt(6) if level == 0 else Pt(4)

        bullet_count += 1

    add_footer(slide)
    return slide

def add_section_slide(prs, section_title):
    """Add section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add large centered text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = title_frame.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide)
    return slide

def clean_markdown(text):
    """Remove markdown formatting"""
    # Remove bold
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'__(.+?)__', r'\1', text)
    # Remove italic
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'_(.+?)_', r'\1', text)
    # Remove inline code
    text = re.sub(r'`(.+?)`', r'\1', text)
    # Remove links
    text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)
    return text.strip()

def parse_markdown_file(md_file):
    """Parse markdown into structured slides"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    slides = []
    current_slide = None
    in_code_block = False
    in_notes = False

    for line in content.split('\n'):
        # Skip notes section
        if '# Presentation Notes' in line or 'Timing Breakdown' in line:
            in_notes = True
        if in_notes:
            continue

        # Handle code blocks
        if line.strip().startswith('```'):
            in_code_block = not in_code_block
            continue
        if in_code_block:
            continue

        # Detect slide markers
        if line.startswith('## SLIDE'):
            # Save previous slide
            if current_slide:
                slides.append(current_slide)

            # Extract slide number and title
            match = re.match(r'## SLIDE (\d+): (.+)', line)
            if match:
                current_slide = {
                    'number': int(match.group(1)),
                    'title': match.group(2).strip(),
                    'bullets': [],
                    'type': 'content'
                }
        elif line.startswith('# ==='):
            # Section header
            if current_slide:
                slides.append(current_slide)
            current_slide = {'type': 'section', 'bullets': []}
        elif line.startswith('#'):
            # Section name
            if current_slide and current_slide['type'] == 'section':
                section_name = line.lstrip('#').strip()
                if 'TRISTAN' in section_name:
                    current_slide['title'] = "CNN Approach\n(Tristan)"
                elif 'BREE' in section_name or 'TRANSFORMER' in section_name:
                    current_slide['title'] = "Transformer Approach\n(Bree)"
                elif 'TEAM' in section_name:
                    current_slide['title'] = "Team Comparison"
                elif 'BACKUP' in section_name:
                    current_slide['title'] = "Backup Slides"
        elif current_slide and line.strip():
            # Add to bullets
            current_slide['bullets'].append(line)

    # Add final slide
    if current_slide:
        slides.append(current_slide)

    return slides

def main():
    md_file = r'C:\CS_4280_Project\term_paper\MIDTERM_SLIDES_APA.md'
    output_file = r'C:\CS_4280_Project\term_paper\MIDTERM_PRESENTATION.pptx'

    print("Parsing markdown...")
    slides_data = parse_markdown_file(md_file)

    print(f"Creating presentation with {len(slides_data)} slides...")
    prs = create_blank_presentation()

    for i, slide_data in enumerate(slides_data):
        if slide_data['type'] == 'section':
            if 'title' in slide_data:
                add_section_slide(prs, slide_data['title'])
        elif slide_data['number'] == 1:
            # Title slide
            add_title_slide(prs, slide_data['title'],
                          ['Josh Manchester', 'Tristan Moffett', 'Brianne Leatherman'])
        else:
            # Content slide
            add_content_slide(prs, slide_data['title'], slide_data['bullets'])

        if (i + 1) % 10 == 0:
            print(f"  Processed {i + 1} slides...")

    print(f"Saving to {output_file}...")
    prs.save(output_file)
    print(f"✓ PowerPoint created successfully!")
    print(f"  Slides: {len(prs.slides)}")

if __name__ == '__main__':
    main()
