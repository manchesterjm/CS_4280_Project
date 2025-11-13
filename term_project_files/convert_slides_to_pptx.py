"""
Convert MIDTERM_SLIDES.md to PowerPoint (.pptx)

Usage:
    python convert_slides_to_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def parse_markdown_slides(md_file):
    """Parse markdown file into slide sections"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by slide markers (## SLIDE or # ===)
    slide_pattern = r'(?:^## SLIDE \d+:.*?$|^# ={3,}$)'
    slides = re.split(slide_pattern, content, flags=re.MULTILINE)

    # Get slide titles
    titles = re.findall(slide_pattern, content, flags=re.MULTILINE)

    parsed_slides = []
    for i, (title, body) in enumerate(zip(titles, slides[1:])):
        # Extract slide number and title
        if title.startswith('## SLIDE'):
            match = re.match(r'## SLIDE (\d+): (.+)', title)
            if match:
                slide_num = match.group(1)
                slide_title = match.group(2).strip()
                parsed_slides.append({
                    'number': slide_num,
                    'title': slide_title,
                    'body': body.strip(),
                    'is_section_header': False
                })
        elif '===' in title:
            # Section header
            parsed_slides.append({
                'number': '',
                'title': '',
                'body': body.strip(),
                'is_section_header': True
            })

    return parsed_slides

def clean_text(text):
    """Clean markdown formatting from text"""
    # Remove code blocks
    text = re.sub(r'```[\s\S]*?```', '', text)
    # Remove bold/italic
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    # Remove inline code
    text = re.sub(r'`(.+?)`', r'\1', text)
    return text.strip()

def add_title_slide(prs, title_text, subtitle_text=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    if subtitle_text:
        subtitle.text = subtitle_text

    return slide

def add_content_slide(prs, title_text, body_text):
    """Add a content slide with title and bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = title_text

    # Add body content
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Parse body text into lines
    lines = body_text.split('\n')

    for line in lines:
        line = line.strip()
        if not line or line.startswith('#') or line.startswith('---'):
            continue

        # Determine indent level
        indent_level = 0
        if line.startswith('  - ') or line.startswith('   - '):
            indent_level = 1
            line = line.lstrip().lstrip('- ').strip()
        elif line.startswith('- '):
            line = line.lstrip('- ').strip()
        elif line.startswith('###'):
            line = line.lstrip('#').strip()
            # Bold for section headers
        elif line.startswith('**') and line.endswith('**'):
            line = line.strip('*').strip()

        if line:
            p = tf.add_paragraph()
            p.text = clean_text(line)
            p.level = indent_level
            p.font.size = Pt(18) if indent_level == 0 else Pt(16)

    return slide

def add_section_header_slide(prs, header_text):
    """Add a section header slide"""
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = header_text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def create_pptx_from_markdown(md_file, output_file):
    """Main conversion function"""
    print(f"Reading {md_file}...")

    # Parse markdown
    slides = parse_markdown_slides(md_file)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print(f"Converting {len(slides)} slides...")

    for i, slide_data in enumerate(slides):
        if slide_data['is_section_header']:
            # Extract section name from body
            section_match = re.search(r'# (.+)', slide_data['body'])
            if section_match:
                section_name = section_match.group(1).strip()
                if 'TRISTAN' in section_name:
                    add_section_header_slide(prs, "CNN Approach\n(Tristan)")
                elif 'BREE' in section_name:
                    add_section_header_slide(prs, "Transformer Approach\n(Bree)")
                elif 'TEAM WRAP-UP' in section_name:
                    add_section_header_slide(prs, "Team Comparison & Wrap-up")
                elif 'BACKUP' in section_name:
                    add_section_header_slide(prs, "Backup Slides")
        else:
            title = slide_data['title']
            body = slide_data['body']

            # Skip presentation notes section
            if 'Presentation Notes' in title or 'Timing Breakdown' in title:
                continue

            # Check if it's a title slide (slide 1)
            if slide_data['number'] == '1':
                # Extract subtitle from body
                subtitle_lines = [line.strip() for line in body.split('\n') if line.strip() and not line.startswith('**')]
                subtitle = '\n'.join(subtitle_lines[:4])
                add_title_slide(prs, title, subtitle)
            else:
                add_content_slide(prs, title, body)

        if (i + 1) % 10 == 0:
            print(f"  Processed {i + 1} slides...")

    # Save presentation
    print(f"Saving to {output_file}...")
    prs.save(output_file)
    print(f"✓ PowerPoint created successfully!")
    print(f"  File: {output_file}")

if __name__ == '__main__':
    md_file = r'C:\CS_4280_Project\term_project_files\MIDTERM_SLIDES.md'
    output_file = r'C:\CS_4280_Project\term_project_files\MIDTERM_PRESENTATION.pptx'

    create_pptx_from_markdown(md_file, output_file)
