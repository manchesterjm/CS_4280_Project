"""
Create the 3 individual paper slides for RNN presentation
These go after Slide 1 (the three papers list)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_content_slide(prs, title):
    """Add blank slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    return slide

def add_text_box(slide, text, left, top, width, height, font_size=18):
    """Add text box to slide"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = text

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)

    return textbox

def add_colored_box(slide, text, left, top, width, height, title, color_rgb, font_size=16):
    """Add colored box with title and text"""
    # Add colored rectangle
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )

    # Set fill color
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color_rgb)

    # Add text
    text_frame = shape.text_frame
    text_frame.word_wrap = True

    # Title
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(font_size + 2)

    # Content
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)

    return shape

# ============================================================================
# SLIDE 2: Machine Learning for Cluster Analysis (Speiser 2020)
# ============================================================================
slide = add_content_slide(prs, "Machine Learning for Cluster Analysis")

# Left box - Key Innovation (Light Gray)
add_colored_box(
    slide,
    "Supervised clustering combined with machine learning for analyzing millions of data points",
    0.5, 1.2, 4, 1.5,
    "Key Innovation",
    (220, 220, 220),  # Light gray
    font_size=14
)

# Right box - Key Takeaway (Dark Blue)
add_colored_box(
    slide,
    "K-means clustering before ML training improves both accuracy and computational speed on large-scale datasets",
    5.5, 1.2, 4, 1.5,
    "Key Takeaway",
    (65, 105, 225),  # Royal blue
    font_size=14
)

# Architecture section
arch_text = """Architecture:
• K-means clustering on extracted features
• Supervised learning on cluster assignments
• Scalable to millions of data points"""

add_text_box(slide, arch_text, 0.5, 3.2, 4.5, 1.5, font_size=16)

# Results section
results_text = """Results:
• Faster training convergence
• Improved classification accuracy
• Handles large-scale datasets efficiently"""

add_text_box(slide, results_text, 5.5, 3.2, 4.5, 1.5, font_size=16)

# Citation at bottom
citation = "Speiser, A., Müller, L. R., Matti, U., Obholzer, N. D., Legant, W. R., Kreshuk, A., ... & Hufnagel, L. (2020). Machine learning for cluster analysis of localization microscopy data. Nature Communications, 11(1), 1493."
add_text_box(slide, citation, 0.5, 6.5, 9, 0.8, font_size=10)

# ============================================================================
# SLIDE 3: LSTM for Time Series Patterns (Vu 2024)
# ============================================================================
slide = add_content_slide(prs, "LSTM for Time Series Patterns")

# Left box - Key Innovation
add_colored_box(
    slide,
    "LSTM networks capture long-term dependencies in noisy environmental time series data",
    0.5, 1.2, 4, 1.5,
    "Key Innovation",
    (220, 220, 220),
    font_size=14
)

# Right box - Key Takeaway
add_colored_box(
    slide,
    "LSTM memory cells handle temporal dependencies that traditional methods miss in sequential data with high noise",
    5.5, 1.2, 4, 1.5,
    "Key Takeaway",
    (65, 105, 225),
    font_size=14
)

# Architecture section
arch_text = """Architecture:
• Stacked LSTM layers for temporal pattern recognition
• Handles irregular sampling and data gaps
• Dropout regularization for noisy data"""

add_text_box(slide, arch_text, 0.5, 3.2, 4.5, 1.5, font_size=16)

# Results section
results_text = """Results:
• Superior performance on noisy time series
• Captures long-term dependencies effectively
• Generalizes to unseen weather patterns"""

add_text_box(slide, results_text, 5.5, 3.2, 4.5, 1.5, font_size=16)

# Citation at bottom
citation = "Vu, M. T., Vo, N. D., Ngo, T. D., Pham, T. D., Huynh, T. T. M., Nguyen, N. T., ... & Ly, H. B. (2024). Harnessing LSTM and XGBoost algorithms for storm prediction. Scientific Reports, 14(1), 11516."
add_text_box(slide, citation, 0.5, 6.5, 9, 0.8, font_size=10)

# ============================================================================
# SLIDE 4: LSTM for Astronomical Photometry (Ding 2024)
# ============================================================================
slide = add_content_slide(prs, "LSTM for Astronomical Photometry")

# Left box - Key Innovation
add_colored_box(
    slide,
    "First application of LSTM to large-scale astronomical photometric flux measurements",
    0.5, 1.2, 4, 1.5,
    "Key Innovation",
    (220, 220, 220),
    font_size=14
)

# Right box - Key Takeaway
add_colored_box(
    slide,
    "LSTM reduced outliers by 33% compared to traditional neural networks on real astronomical survey data",
    5.5, 1.2, 4, 1.5,
    "Key Takeaway",
    (65, 105, 225),
    font_size=14
)

# Architecture section
arch_text = """Architecture:
• LSTM for photometric time series
• Handles large survey datasets
• Optimized for astronomical flux measurements"""

add_text_box(slide, arch_text, 0.5, 3.2, 4.5, 1.5, font_size=16)

# Results section
results_text = """Results:
• 33% fewer outliers than MLPs
• Better handling of temporal photometric data
• Validated on real survey observations"""

add_text_box(slide, results_text, 5.5, 3.2, 4.5, 1.5, font_size=16)

# Citation at bottom
citation = "Ding, Y., Ji, K., Xiao, M., Zheng, X., Chen, Y., Liang, J., ... & Qi, Y. (2024). Photometric redshift estimation for CSST survey with LSTM neural networks. Monthly Notices of the Royal Astronomical Society, 535(2), 1844-1858."
add_text_box(slide, citation, 0.5, 6.5, 9, 0.8, font_size=10)

# ============================================================================
# Save presentation
# ============================================================================
output_path = Path(r"C:\CS_4280_Project\term_project_files\PAPER_SLIDES_TO_ADD.pptx")
prs.save(str(output_path))

print(f"PowerPoint created: {output_path}")
print(f"3 slides ready to copy into your presentation")
print(f"These go AFTER your first slide (the 3 papers list)")
