"""
Create RNN Midterm Presentation PowerPoint
Generates complete 9-slide presentation with embedded images
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title):
    """Add blank slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    return slide

def add_text(slide, text, left, top, width, height, font_size=14):
    """Add text box to slide"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for line in text.split('\n'):
        if line.strip():
            p = text_frame.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(font_size)
            p.level = 0

    # Remove first empty paragraph
    if text_frame.paragraphs[0].text == '':
        text_frame.clear()
        for line in text.split('\n'):
            if line.strip():
                p = text_frame.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(font_size)

    return textbox

def add_image(slide, image_path, left, top, width=None, height=None):
    """Add image to slide"""
    if Path(image_path).exists():
        if width and height:
            slide.shapes.add_picture(str(image_path), Inches(left), Inches(top),
                                    width=Inches(width), height=Inches(height))
        elif width:
            slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), height=Inches(height))
        else:
            slide.shapes.add_picture(str(image_path), Inches(left), Inches(top))
        return True
    return False

# Paths
base_path = Path(r"C:\CS_4280_Project\term_project_files")
fig_path = base_path / "materials" / "figures" / "rnn_slides"

# ============================================================================
# SLIDE 1: Related Work - Three NEW Papers
# ============================================================================
slide = add_content_slide(prs, "Related Work - Three NEW Papers")

text = """Speiser, A., Müller, L. R., Matti, U., Obholzer, N. D., Legant, W. R., Kreshuk, A., ... & Hufnagel, L. (2020). Machine learning for cluster analysis of localization microscopy data. Nature Communications, 11(1), 1493. (H5 index: 399)

Vu, M. T., Vo, N. D., Ngo, T. D., Pham, T. D., Huynh, T. T. M., Nguyen, N. T., ... & Ly, H. B. (2024). Harnessing LSTM and XGBoost algorithms for storm prediction. Scientific Reports, 14(1), 11516. (H5 index: 234)

Ding, Y., Ji, K., Xiao, M., Zheng, X., Chen, Y., Liang, J., ... & Qi, Y. (2024). Photometric redshift estimation for CSST survey with LSTM neural networks. Monthly Notices of the Royal Astronomical Society, 535(2), 1844-1858. (H5 index: 151)"""

add_text(slide, text, 0.5, 1.2, 9, 5.5, font_size=12)

# ============================================================================
# SLIDE 2: Why BiLSTM + Clustering?
# ============================================================================
slide = add_content_slide(prs, "Why BiLSTM + Clustering?")

text = """From the Papers:
• Clustering improves ML on large datasets (Speiser 2020)
• LSTM captures long-term dependencies (Vu 2024)
• LSTM reduces outliers in astronomy (Ding 2024)

My Approach:
BiLSTM + K-means clustering on BLS features"""

add_text(slide, text, 0.5, 1.5, 9, 4, font_size=20)

# ============================================================================
# SLIDE 3: Methodology
# ============================================================================
slide = add_content_slide(prs, "Methodology")

text = """Data: 655 windows (150 planets, 505 non-planets)

Features: Period, depth, duration, BLS power

Architecture:
• K-means (5 clusters)
• 4-layer BiLSTM (256 hidden, bidirectional)
• Cluster embeddings (32-dim)"""

add_text(slide, text, 0.5, 1.2, 9, 2.5, font_size=16)

# Add preprocessing diagram
add_image(slide, fig_path / "preprocessing_pipeline.png", 0.5, 4, width=9)

# ============================================================================
# SLIDE 4: BiLSTM Architecture
# ============================================================================
slide = add_content_slide(prs, "BiLSTM Architecture")

text = """Input (2048) → BiLSTM Layers → Cluster Embedding → FC Layers → Output

~2.1M parameters"""

add_text(slide, text, 0.5, 1.2, 9, 1, font_size=16)

# Add architecture diagram
add_image(slide, fig_path / "bilstm_architecture.png", 1, 2.5, width=8)

# ============================================================================
# SLIDE 5: Results
# ============================================================================
slide = add_content_slide(prs, "Results")

text = """AUC: 75.72%
Recall: 88.67% | Precision: 38.27%

Tested on 100 confirmed exoplanet systems"""

add_text(slide, text, 0.5, 1.2, 4, 1.5, font_size=16)

# Add metrics bar chart (left)
add_image(slide, fig_path / "metrics_bar_chart.png", 0.5, 2.5, width=4.5)

# Add confusion matrix (right)
add_image(slide, fig_path / "confusion_matrix.png", 5.2, 2.5, width=4.5)

# ============================================================================
# SLIDE 6: Learning from Failure
# ============================================================================
slide = add_content_slide(prs, "Learning from Failure")

text = """100 Planets Only:
❌ Predicted everything as planet
❌ 100% false positives

Solution: Add Non-Planets
✓ 100 planets + 300 non-planets → AUC: 0.69
✓ Real TESS: Identified TIC 307210830 (prob: 0.5959)

But: Imbalanced Data Problem
❌ 150 planets vs 505 non-planets (23% positive)
❌ High recall (88.67%) but low precision (38.27%)
❌ Too many false positives → Model biased toward negatives"""

add_text(slide, text, 0.5, 1.5, 9, 5, font_size=16)

# ============================================================================
# SLIDE 7: Optuna Optimization
# ============================================================================
slide = add_content_slide(prs, "Optuna Optimization")

text = """Parameter changes:
• Layers: 3 → 4
• Batch size: 64 → 128
• Learning rate: 1e-4 → 2.25e-4

Result: AUC 0.69 → 0.76 (+9%)"""

add_text(slide, text, 0.5, 1.2, 4.5, 2, font_size=16)

# Add progression chart
add_image(slide, fig_path / "model_progression.png", 5, 1.5, width=4.5)

# ============================================================================
# SLIDE 8: Demo
# ============================================================================
slide = add_content_slide(prs, "Demo")

text = """Model identifying TIC 307210830
(L 98-59 multi-planet system)

[Demo video shown during presentation]"""

add_text(slide, text, 2, 3, 6, 2, font_size=20)

# Note: Video embedding in PowerPoint is complex, user can add manually

# ============================================================================
# SLIDE 9: What's Next?
# ============================================================================
slide = add_content_slide(prs, "What's Next?")

text = """Balanced Data Attempt Failed
• Tried 50/50 balanced synthetic data (200 planets + 200 non-planets)
• AUC dropped to 0.45 on real TESS data (domain shift)

Solution: Hybrid Training
• Mix real TESS + synthetic (90/10 ratio)
• Better balance than 23% but maintains domain fidelity

Cross-Mission Testing
Train on TESS → Test on Kepler to verify generalization"""

add_text(slide, text, 0.5, 1.5, 9, 5, font_size=16)

# ============================================================================
# Save presentation
# ============================================================================
output_path = base_path / "RNN_MIDTERM_PRESENTATION.pptx"
prs.save(str(output_path))

print(f"PowerPoint created: {output_path}")
print(f"9 slides with embedded images")
print(f"Ready to present!")
